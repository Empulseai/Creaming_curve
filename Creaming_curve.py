import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import streamlit as st
from openpyxl import Workbook
from openpyxl.styles import Alignment,PatternFill,Font
import xlsxwriter
import win32com.client as win32
import pythoncom
import random
from openpyxl.drawing.image import Image
import os
import time
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from datetime import date
from datetime import datetime
import matplotlib.dates as mdates
import spacy
 
st.markdown(
"""
    <style>
        body {
            background-color: #f4f4f4;
            font-family: 'Helvetica Neue', sans-serif;
        }
        .stApp {
            background-color: #f4f4f4;
        }
        h1 {
            color: #0078D7;
            text-align: center;
            font-size: 36px;
            font-family: 'Georgia', serif;
        }
        h4 {
            text-align: center;
            color: #444;
            font-family: 'Verdana', sans-serif;
        }
        .css-18e3th9 {
            padding-top: 2rem;
        }
        .stDataFrame {
            background-color: white;
            border-radius: 10px;
            padding: 10px;
        }
    </style>
    """,
    unsafe_allow_html=True
)
 
 
st.title("EI MATRIX TOOL")
file_upload1=st.file_uploader("📂 Upload the Excel file for EI matrix",type=["xlsx"],accept_multiple_files=False)
if file_upload1 is not None:
    df=pd.read_excel(file_upload1,sheet_name="Consolidated Ideas")
    df["Count of the category"]=df["Category"].map(df["Category"].value_counts())
    df=df.sort_values(by=["Count of the category","Category"],ascending=False).reset_index(drop=True)
    df["Sl No"]=range(1,len(df)+1)
    df["Effort1"]=np.where(df["Effort"]=="High",1.5,0.5)
    df["Impact1"]=np.where(df["Impact"]=="High",1.5,0.5)
    pythoncom.CoInitialize()
    excel = win32.Dispatch('Excel.Application')
    excel.Visible = True
    workbook = excel.Workbooks.Add()
    worksheet = workbook.Worksheets(1)
    worksheet.Cells(1, 1).Value = "Pain Points and Improvement Oppertunity"
    worksheet.Cells(1, 1).Font.Size=20
    worksheet.Cells(1, 1).Font.Bold=True
    worksheet.Cells(1, 22).Value = "Consolidation and Cluster OPPertunity"
    worksheet.Cells(1, 22).Font.Size=20
    worksheet.Cells(1,22).Font.Bold=True
 
    start_top = 50
    start_left = 0
    rectangle_width = 100
    rectangle_height = 70
    vertical_spacing = 80
    horizontal_spacing = 110
    total_rectangle=len(df)
    for idx, row in df.iterrows():
        text_from_df = row['Category']
        additional_text = row['Ideas']
        sl_no=row['Sl No']
        combined_text = f"{sl_no}\n{text_from_df}\n\n{additional_text}"
   
 
        row_position = idx // 5  
        col_position = idx % 5
       
        top = start_top + row_position * vertical_spacing  
        left = start_left + col_position * horizontal_spacing
       
        shape = worksheet.Shapes.AddTextbox(1, left, top, rectangle_width, rectangle_height)  # Parameters: Type, Left, Top, Width, Height
        shape.Fill.ForeColor.RGB = 0xFFFFCD  # Set the color (Tomato red)
        shape.Line.Weight = 2  # Set the border thickness
        shape.TextFrame2.TextRange.Text = combined_text  # Add text to the shape
 
    # Optionally, you can format the text (e.g., font size, bold)
        shape.TextFrame2.TextRange.Font.Size = 8  # Font size
        shape.TextFrame2.TextRange.Font.Bold = True
    start_top = 50
    start_left = 1000
    rectangle_width = 100
    rectangle_height = 70
    vertical_spacing = 80
    horizontal_spacing = 110
    cluster_vertical_spacing = 300  # Vertical space between different clusters
    cluster_horizontal_spacing = 250  # Horizontal space between clusters
    category_top_offset = 0
    grouped = df.groupby('Category')
 
 
    def generate_random_color():
        r = random.randint(128, 255)  # Red component
        g = random.randint(128, 255)  # Green component
        b = random.randint(128, 255)
        return (r,g,b)
    unique_categories = df['Category'].unique()
    color_mapping = {category: generate_random_color() for category in unique_categories}
   
 
    category_positions = {category: (0, 1000) for category in unique_categories}
    category_counts = {category: 0 for category in unique_categories}
    color_groups = {}
    for idx, row in df.iterrows():
    # Get the text for the rectangle from the DataFrame
        text_from_df = row['Category']
        additional_text = row['Ideas']
        sl_no=row['Sl No']
   
    # Combine the texts (you can adjust the separator as needed)
        # combined_text = text_from_df + "\n\n" + additional_text
       
        combined_text = f"{sl_no}\n{text_from_df}\n\n{additional_text}"
       
 
    # Calculate the position of each rectangle
        row_position = idx // 5  # Integer division to determine row (0 for the first row, 1 for the second row, etc.)
        col_position = idx % 5
       
        top = start_top + row_position * vertical_spacing  # Move down after every 5 rectangles (1 row)
        left = start_left + col_position * horizontal_spacing
       
       
       
        shape = worksheet.Shapes.AddTextbox(1,left , top, rectangle_width, rectangle_height)  # Parameters: Type, Left, Top, Width, Height
        fill_color = color_mapping.get(text_from_df,0xFFFFFF)
        rgb_value = (fill_color[0] << 16) + (fill_color[1] << 8) + fill_color[2]
        shape.Fill.ForeColor.RGB = rgb_value# Set the color (Tomato red)
        shape.Line.Weight = 1  # Set the border thickness
        shape.TextFrame2.TextRange.Text = combined_text
        # Add text to the shape
 
    # Optionally, you can format the text (e.g., font size, bold)
       
        shape.TextFrame2.TextRange.Font.Size = 8  # Font size
 
        shape.TextFrame2.TextRange.Font.Bold = True
   
    filtered_df = df[(df['Effort1'] == 0.5) & (df['Impact1'] == 1.5)]
    df1=pd.DataFrame(filtered_df)
    df1=df1.reset_index(drop=True)
    df1["Sl No"]=range(1,len(df1)+1)
    start_top = 300
    start_left = 2000
    rectangle_width = 100
    rectangle_height = 70
    vertical_spacing = 80
    horizontal_spacing = 110
    for idx, row in df1.iterrows():
        # Get the text for the rectangle from the DataFrame
        text_from_df = row['Category']
        additional_text = row['Ideas']
        sl_no=row['Sl No']
   
       
        combined_text = f"{sl_no}\n{text_from_df}\n\n{additional_text}"
       
 
    # Calculate the position of each rectangle
        row_position = idx // 8  # Integer division to determine row (0 for the first row, 1 for the second row, etc.)
        col_position = idx % 8
       
        top = start_top + row_position * vertical_spacing  # Move down after every 5 rectangles (1 row)
        left = start_left + col_position * horizontal_spacing
       
       
       
        shape = worksheet.Shapes.AddTextbox(1,left , top, rectangle_width, rectangle_height)  # Parameters: Type, Left, Top, Width, Height
        fill_color = color_mapping.get(text_from_df,0xFFFFFF)
        rgb_value = (fill_color[0] << 16) + (fill_color[1] << 8) + fill_color[2]
        shape.Fill.ForeColor.RGB = rgb_value # Set the color (Tomato red)
        shape.Line.Weight = 1  # Set the border thickness
        shape.TextFrame2.TextRange.Text = combined_text
        # Add text to the shape
 
    # Optionally, you can format the text (e.g., font size, bold)
       
        shape.TextFrame2.TextRange.Font.Size = 8  # Font size
 
        shape.TextFrame2.TextRange.Font.Bold = True
    filtered_df_1 = df[(df['Effort1'] == 0.5) & (df['Impact1'] == 0.5)]
    df2=pd.DataFrame(filtered_df_1)
    df2=df2.reset_index(drop=True)
    df2["Sl No"]=range(1,len(df2)+1)
    start_top = 1200
    start_left = 2000
    rectangle_width = 100
    rectangle_height = 70
    vertical_spacing = 80
    horizontal_spacing = 110
    for idx, row in df2.iterrows():
        # Get the text for the rectangle from the DataFrame
        text_from_df = row['Category']
        additional_text = row['Ideas']
        sl_no=row['Sl No']
   
       
        combined_text = f"{sl_no}\n{text_from_df}\n\n{additional_text}"
       
 
    # Calculate the position of each rectangle
        row_position = idx // 8  # Integer division to determine row (0 for the first row, 1 for the second row, etc.)
        col_position = idx % 8
       
        top = start_top + row_position * vertical_spacing  # Move down after every 5 rectangles (1 row)
        left = start_left + col_position * horizontal_spacing
       
       
       
        shape = worksheet.Shapes.AddTextbox(1,left , top, rectangle_width, rectangle_height)  # Parameters: Type, Left, Top, Width, Height
        fill_color = color_mapping.get(text_from_df,0xFFFFFF)
        rgb_value = (fill_color[0] << 16) + (fill_color[1] << 8) + fill_color[2]
        shape.Fill.ForeColor.RGB = rgb_value  # Set the color (Tomato red)
        shape.Line.Weight = 1  # Set the border thickness
        shape.TextFrame2.TextRange.Text = combined_text
        # Add text to the shape
 
    # Optionally, you can format the text (e.g., font size, bold)
       
        shape.TextFrame2.TextRange.Font.Size = 8  # Font size
 
        shape.TextFrame2.TextRange.Font.Bold = True
    filtered_df_2 = df[(df['Effort1'] == 1.5) & (df['Impact1'] == 1.5)]
    df3=pd.DataFrame(filtered_df_2)
    df3=df3.reset_index(drop=True)
    df3["Sl No"]=range(1,len(df3)+1)
    start_top = 300
    start_left = 3000
    rectangle_width = 100
    rectangle_height = 70
    vertical_spacing = 80
    horizontal_spacing = 110
    for idx, row in df3.iterrows():
        # Get the text for the rectangle from the DataFrame
        text_from_df = row['Category']
        additional_text = row['Ideas']
        sl_no=row['Sl No']
   
       
        combined_text = f"{sl_no}\n{text_from_df}\n\n{additional_text}"
       
 
    # Calculate the position of each rectangle
        row_position = idx // 8  # Integer division to determine row (0 for the first row, 1 for the second row, etc.)
        col_position = idx % 8
       
        top = start_top + row_position * vertical_spacing  # Move down after every 5 rectangles (1 row)
        left = start_left + col_position * horizontal_spacing
       
       
       
        shape = worksheet.Shapes.AddTextbox(1,left , top, rectangle_width, rectangle_height)  # Parameters: Type, Left, Top, Width, Height
        fill_color = color_mapping.get(text_from_df,0xFFFFFF)
        rgb_value = (fill_color[0] << 16) + (fill_color[1] << 8) + fill_color[2]
        shape.Fill.ForeColor.RGB = rgb_value  # Set the color (Tomato red)
        shape.Line.Weight = 1  # Set the border thickness
        shape.TextFrame2.TextRange.Text = combined_text
        # Add text to the shape
 
    # Optionally, you can format the text (e.g., font size, bold)
       
        shape.TextFrame2.TextRange.Font.Size = 8  # Font size
 
        shape.TextFrame2.TextRange.Font.Bold = True
    filtered_df_3 = df[(df['Effort1'] == 1.5) & (df['Impact1'] == 0.5)]
    df4=pd.DataFrame(filtered_df_3)
    df4=df4.reset_index(drop=True)
    df4["Sl No"]=range(1,len(df4)+1)
    start_top = 1200
    start_left = 3000
    rectangle_width = 100
    rectangle_height = 70
    vertical_spacing = 80
    horizontal_spacing = 110
    for idx, row in df4.iterrows():
        # Get the text for the rectangle from the DataFrame
        text_from_df = row['Category']
        additional_text = row['Ideas']
        sl_no=row['Sl No']
   
       
        combined_text = f"{sl_no}\n{text_from_df}\n\n{additional_text}"
       
 
    # Calculate the position of each rectangle
        row_position = idx // 8  # Integer division to determine row (0 for the first row, 1 for the second row, etc.)
        col_position = idx % 8
       
        top = start_top + row_position * vertical_spacing  # Move down after every 5 rectangles (1 row)
        left = start_left + col_position * horizontal_spacing
       
       
       
        shape = worksheet.Shapes.AddTextbox(1,left , top, rectangle_width, rectangle_height)  # Parameters: Type, Left, Top, Width, Height
        fill_color = color_mapping.get(text_from_df,0xFFFFFF)
        rgb_value = (fill_color[0] << 16) + (fill_color[1] << 8) + fill_color[2]
        shape.Fill.ForeColor.RGB = rgb_value  # Set the color (Tomato red)
        shape.Line.Weight = 1  # Set the border thickness
        shape.TextFrame2.TextRange.Text = combined_text
        # Add text to the shape
 
    # Optionally, you can format the text (e.g., font size, bold)
       
        shape.TextFrame2.TextRange.Font.Size = 8  # Font size
 
        shape.TextFrame2.TextRange.Font.Bold = True
    data={
        "Category":["Quick Wins","Fill Ins","Major Project","Time Wasters"],
        "X_values":[0.8,0.8,1.8,1.8],
        "Y_values":[1.2,0.2,1.2,0.2]
    }
    data=pd.DataFrame(data)
    plt.figure(figsize=(10, 6))
    plt.scatter(data["X_values"], data["Y_values"], marker=' ')
    for i, txt in enumerate(data['Category']):
        plt.annotate(txt, (data['X_values'][i], data['Y_values'][i]), fontsize=9, ha='right')
    light_gray = (242/255, 242/255, 242/255)
    # Set the title and labels
    ax = plt.gca()
    ax.set_facecolor(light_gray)
    plt.title('EI MATRIX')
    ax.set_xlabel("EFFORT")
    ax.set_ylabel("IMPACT")
    ax.set_xlim(0,2)
    ax.set_xticks(np.arange(0,2,0.5))
    ax.set_xticklabels([" ","Low"," ","High"])
    ax.set_ylim(0,2)
    ax.set_yticks(np.arange(0,2,0.5))
    ax.set_yticklabels([" ","Low"," ","High"])
    ax.spines['top'].set_color('white')
    ax.spines['top'].set_linewidth(0)
    ax.spines['right'].set_color('white')
    ax.spines['right'].set_linewidth(0)
    ax.spines['left'].set_color('black')
    ax.spines['left'].set_linewidth(2)
    ax.spines['bottom'].set_color('black')
    ax.spines['bottom'].set_linewidth(2)
 
    # Show the grid
    for x in [0, 1]:  # x-ticks without 0.5 and 1.5
        ax.axvline(x=x, color='black', linewidth=1.5,linestyle="--")
 
    for y in [0, 1]:  # y-ticks without 0.5 and 1.5
        ax.axhline(y=y, color='black', linewidth=1.5,linestyle="--")
 
    # Show the plot
    image_path = 'scatter_plot.png'
    plt.savefig(image_path)  # Save as PNG file
    plt.close()
    picture=worksheet.Shapes.AddPicture(
        Filename=os.path.abspath(image_path),  # Use absolute path for the image
        LinkToFile=False,
        SaveWithDocument=True,
        Left=1650,  # Position from the left
        Top=0,   # Position from the top
        Width=2500,  # Use default width
        Height=2200 # Use default height
    )
    time.sleep(1)
    picture.ZOrder(1)
 
    output_excel = "output_with_shape.xlsx"
    workbook.SaveAs(output_excel)
    st.dataframe(df1)
    st.data_editor(df)
 
# Close the workbook and Excel
    workbook.Close()
   
    excel.Quit()
    pythoncom.CoUninitialize()
 
 
 
################## CREAMING CURVE ###########################
 
 
 
st.title("Creaming Curve Analyzer")
 
 
file=st.file_uploader("📂 Upload Excel file for Creaming curve ", type=["xlsx"], help="Ensure the file contains relevant cost and savings data")
 
 
if file is not None:
    df = pd.read_excel(file)
 
    if 'Cost $' not in df.columns or 'Annual Savings $ K' not in df.columns:
        st.error("⚠️ Data is incomplete! Please enter 'Cost $' and 'Annual Savings $ K' values.")
        df['Cost $'] = 0
        df['Annual Savings $ K'] = 0  
    df = st.data_editor(df, num_rows="dynamic")
 
    df['Savings ratio'] = df['Annual Savings $ K'] / df['Cost $']
    df = df.sort_values(by='Savings ratio', ascending=False)
    df['Cumulative cost'] = df['Cost $'].cumsum()
    df['Cumulative Savings'] = df['Annual Savings $ K'].cumsum()
 
 
    st.markdown("Updated DataFrame:")
    st.dataframe(df.style.format({"Cost $": "${:,.2f}", "Annual Savings $ K": "${:,.2f}"}))
 
    fig, ax = plt.subplots(figsize=(12, 6))
    fig.patch.set_facecolor('#f4f4f4')
 
    budget = st.number_input("💰 Enter your budget ($ K):", min_value=0.0, step=100.0)
 
 
    if budget > 0:
        within_budget = df['Cumulative cost'] <= budget
        ax.scatter(df.loc[within_budget, 'Cumulative cost'], df.loc[within_budget, 'Cumulative Savings'],
                   color='green', edgecolors='black', s=100, alpha=0.8, label='Projects within Budget')
        ax.scatter(df.loc[~within_budget, 'Cumulative cost'], df.loc[~within_budget, 'Cumulative Savings'],
                   color='red', edgecolors='black', s=100, alpha=0.8, label='Projects outside Budget')
        ax.axvline(x=budget, color='red', linestyle='--', linewidth=2, label=f'Budget: ${budget:,.0f} K')
    else:
        ax.scatter(df['Cumulative cost'], df['Cumulative Savings'],
                   color='blue', edgecolors='black', s=100, alpha=0.8, label='Projects')
         
    ax.set_title('Cumulative Cost vs. Cumulative Savings', fontsize=14, fontweight='bold', color='#333', fontname='Courier New')
    ax.set_xlabel('Cumulative Cost ($ K)', fontsize=10, fontweight='bold', color='#444', fontname='Tahoma')
    ax.set_ylabel('Cumulative Savings ($ K)', fontsize=10, fontweight='bold', color='#444', fontname='Tahoma')    
 
    ax.grid(True, linestyle='--', alpha=0.5, color='gray')
 
    x_labels = [f"{name}  (${cost:,.0f})" for name, cost in zip(df['Project Summary Name'], df['Cumulative cost'])]
    plt.xticks(df['Cumulative cost'], x_labels, rotation=90, ha='center', fontsize=8, color='#222', fontname='Arial')
 
    ax.legend(loc='upper left', frameon=True, facecolor='white', edgecolor='black')
    st.pyplot(fig)
 
 
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"Creaming Curve                                 with Budget ${budget:,.0f}"
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight')
    img_stream.seek(0)
    slide.shapes.add_picture(img_stream, Inches(1), Inches(1), width=Inches(8))
    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
 
 
    st.download_button("📥 Download PowerPoint", ppt_stream, "creaming_curve_presentation.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
 
    excel_stream = BytesIO()
    with pd.ExcelWriter(excel_stream, engine='openpyxl') as writer:
     df.to_excel(writer, index=False, sheet_name='DataFrame')
     worksheet = writer.sheets['DataFrame']
   
     fill_green = PatternFill(start_color="90EE90", end_color="90EE90", fill_type="solid")  # Light green
     fill_red = PatternFill(start_color="FF9999", end_color="FF9999", fill_type="solid")  # Light red
     font_green = Font(color="006400", bold=False)  # White bold font for headers
     font_red = Font(color="8B0000", bold=False)  # White bold font for headers
   
     fill_blue = PatternFill(start_color="0000FF", end_color="0000FF", fill_type="solid")  # Blue header
     font_white_bold = Font(color="FFFFFF", bold=True)  # White bold font for headers
 
     for cell in worksheet[1]:  # First row (header row)
        cell.fill = fill_blue
        cell.font = font_white_bold
   
     for row in range(2, len(df) + 2):  
        is_within_budget = df.iloc[row - 2]['Cumulative cost'] <= budget
        fill_color = fill_green if is_within_budget else fill_red
        fill_font=font_green if is_within_budget else font_red
       
        for col in worksheet.iter_cols(min_row=row, max_row=row, min_col=1, max_col=len(df.columns)):
            for cell in col:
                # cell.fill = fill_color
                cell.font= fill_font
 
    excel_stream.seek(0)
   
    st.download_button("📥 Download Excel", excel_stream, "creaming_curve_data.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
   
    st.success("✅ Analysis Complete!")
 